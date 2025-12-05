from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Alias for compatibility
RgbColor = RGBColor

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    sub_frame = sub_box.text_frame
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = subtitle
    sub_para.font.size = Pt(24)
    sub_para.font.color.rgb = RgbColor(100, 100, 100)
    sub_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0, 51, 102)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = content_frame.paragraphs[0]
        else:
            para = content_frame.add_paragraph()
        para.text = "• " + bullet
        para.font.size = Pt(20)
        para.space_after = Pt(12)

def add_diagram_slide(prs, title, content_text):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0, 51, 102)

    # Content box
    content_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.2), Inches(9.4), Inches(5.6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    para = content_frame.paragraphs[0]
    para.text = content_text
    para.font.size = Pt(14)
    para.font.name = "Consolas"

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(0, 51, 102)

    # Table
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.3), Inches(1.3), Inches(9.4), Inches(0.5)).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(12)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
add_title_slide(prs, "Document Hub POC", "Configurable Data Extraction & Dynamic Document Matching")

# Slide 2: Overview
add_content_slide(prs, "System Overview", [
    "Document Hub provides intelligent document retrieval based on customer/account context",
    "Supports multiple document sharing scopes: ALL, ACCOUNT, CUSTOMER, CUSTOM_RULES",
    "CUSTOM_RULES enables dynamic document matching via configurable API chains",
    "Two main matching strategies: Reference Key Matching & Conditional Matching",
    "Uses PostgreSQL with R2DBC for reactive database access"
])

# Slide 3: Architecture
add_content_slide(prs, "Key Components", [
    "DocumentEnquiryController - REST API endpoint for document queries",
    "DocumentEnquiryService - Core business logic for document retrieval",
    "ConfigurableDataExtractionService - Executes API chains, extracts fields",
    "MasterTemplateDefinition - Template configuration with data_extraction_config",
    "StorageIndex - Document storage with reference keys for matching",
    "MockApiController - Simulates external APIs for testing"
])

# Slide 4: Data Extraction Config
config_example = '''data_extraction_config (JSONB):
{
  "requiredFields": ["disclosureCode"],
  "fieldSources": {
    "disclosureCode": {
      "sourceApi": "pricingApi",
      "extractionPath": "$.disclosureCode",
      "requiredInputs": ["pricingId"]
    }
  },
  "dataSources": {
    "accountArrangementsApi": {
      "endpoint": { "url": "http://.../accounts/${accountId}/arrangements" },
      "providesFields": ["pricingId"]
    },
    "pricingApi": {
      "endpoint": { "url": "http://.../pricing/${pricingId}" },
      "providesFields": ["disclosureCode"]
    }
  },
  "documentMatching": {
    "matchBy": "reference_key",
    "referenceKeyField": "disclosureCode"
  }
}'''
add_diagram_slide(prs, "Data Extraction Configuration", config_example)

# Slide 5: Reference Key Matching Flow
add_content_slide(prs, "Reference Key Matching Flow", [
    "1. Request received with customerId and accountId",
    "2. Load all active templates from master_template_definition",
    "3. For CUSTOM_RULES templates, execute API chain:",
    "   - Call accountArrangementsApi → extract pricingId",
    "   - Call pricingApi with pricingId → extract disclosureCode",
    "4. Use disclosureCode as reference_key to query storage_index",
    "5. Return matching document(s)",
    "",
    "Example: Account with pricingId 'PR001' → disclosureCode 'D164'",
    "         → Returns 'Credit_Card_Terms_D164_v1.pdf'"
])

# Slide 6: Conditional Matching
conditional_example = '''documentMatching (for TIERED_CARD_AGREEMENT):
{
  "matchBy": "conditional",
  "referenceKeyType": "CREDIT_TIER",
  "conditions": [
    { "field": "creditLimit", "operator": ">=", "value": 50000,
      "referenceKey": "TIER-PLATINUM" },
    { "field": "creditLimit", "operator": ">=", "value": 25000,
      "referenceKey": "TIER-GOLD" },
    { "field": "creditLimit", "operator": "<", "value": 25000,
      "referenceKey": "TIER-STANDARD" }
  ]
}

Supported Operators: >=, >, <=, <, ==, !=
Field Types: Number, String, Boolean'''
add_diagram_slide(prs, "Conditional Matching Configuration", conditional_example)

# Slide 7: Conditional Matching Flow
add_content_slide(prs, "Conditional Matching Flow", [
    "1. Request received with accountId",
    "2. Call creditInfoApi to get creditLimit for the account",
    "3. Evaluate conditions in order (first match wins):",
    "   - creditLimit >= 50000 → TIER-PLATINUM",
    "   - creditLimit >= 25000 → TIER-GOLD",
    "   - creditLimit < 25000 → TIER-STANDARD",
    "4. Query storage_index with matched reference_key",
    "5. Return the tier-appropriate document",
    "",
    "Error Handling: If API fails or no condition matches,",
    "template is SKIPPED (returns empty) - no fallthrough"
])

# Slide 8: Test Accounts Table
add_table_slide(prs, "Test Accounts - Conditional Matching",
    ["Account ID", "Credit Limit", "Tier", "Document"],
    [
        ["aaaa...0001", "$75,000", "PLATINUM", "Platinum_Card_Agreement.pdf"],
        ["aaaa...0002", "$35,000", "GOLD", "Gold_Card_Agreement.pdf"],
        ["aaaa...0003", "$15,000", "STANDARD", "Standard_Card_Agreement.pdf"]
    ])

# Slide 9: Disclosure Code Table
add_table_slide(prs, "Test Accounts - Disclosure Code Matching",
    ["Account", "Disclosure Code", "Document"],
    [
        ["Account 1 (VIP)", "D164", "Credit_Card_Terms_D164_v1.pdf"],
        ["Account 2", "D166", "Premium_Credit_Card_Terms_D166_v1.pdf"],
        ["Account 3", "D165", "Credit_Card_Terms_D165_v1.pdf"]
    ])

# Slide 10: API Flow Diagram
flow_diagram = '''
┌─────────────────────────────────────────────────────────────────────┐
│                     Document Enquiry Flow                           │
└─────────────────────────────────────────────────────────────────────┘

   Client Request
        │
        ▼
┌───────────────────┐
│ POST /documents-  │
│     enquiry       │
└─────────┬─────────┘
          │
          ▼
┌───────────────────┐     ┌─────────────────────┐
│ Load Active       │────▶│ master_template_    │
│ Templates         │     │ definition          │
└─────────┬─────────┘     └─────────────────────┘
          │
          ▼
┌───────────────────┐     ┌─────────────────────┐
│ For CUSTOM_RULES: │────▶│ External APIs       │
│ Execute API Chain │     │ (Mock/Real)         │
└─────────┬─────────┘     └─────────────────────┘
          │
          ▼
┌───────────────────┐     ┌─────────────────────┐
│ Match Documents   │────▶│ storage_index       │
│ by Reference Key  │     │                     │
└─────────┬─────────┘     └─────────────────────┘
          │
          ▼
┌───────────────────┐
│ Return Document   │
│ List Response     │
└───────────────────┘
'''
add_diagram_slide(prs, "Document Enquiry Flow", flow_diagram)

# Slide 11: Key Features
add_content_slide(prs, "Key Features Implemented", [
    "Dynamic API Chain Execution - sequential calls with dependency resolution",
    "JSONPath Field Extraction - $.field.nested.path support",
    "Variable Substitution - ${accountId}, ${pricingId} in URLs",
    "Conditional Document Matching - comparison operators for numeric/string fields",
    "Error Handling - graceful skip on API failures (no fallthrough)",
    "Array Unwrapping - single-element arrays auto-unwrapped",
    "Configurable Timeouts - per-API timeout settings"
])

# Slide 12: Required Headers
add_content_slide(prs, "API Request Requirements", [
    "Required Headers:",
    "  • Content-Type: application/json",
    "  • X-version: 1",
    "  • X-correlation-id: <unique-request-id>",
    "  • X-requestor-id: <UUID>",
    "  • X-requestor-type: CUSTOMER | BANKER",
    "",
    "Request Body:",
    "  { \"customerId\": \"<uuid>\", \"accountId\": [\"<uuid>\"] }"
])

# Slide 13: Test Scripts
add_content_slide(prs, "Test Scripts Available", [
    "Location: fresh-doc-hub-poc/test-scripts/",
    "",
    "• restart-app.bat - Kill and restart Spring Boot application",
    "• run-all-tests.bat - Execute all document enquiry tests",
    "• test-conditional-platinum.bat - Test Platinum tier ($75k)",
    "• test-conditional-gold.bat - Test Gold tier ($35k)",
    "• test-conditional-standard.bat - Test Standard tier ($15k)",
    "• test-disclosure-code-d164.bat - Test D164 disclosure matching",
    "• test-disclosure-code-d166.bat - Test D166 disclosure matching",
    "• API_REFERENCE.md - Complete curl command reference"
])

# Slide 14: Summary
add_content_slide(prs, "Summary", [
    "Document Hub POC demonstrates flexible document retrieval:",
    "",
    "✓ Reference Key Matching - API chain extracts key, matches documents",
    "✓ Conditional Matching - Evaluate rules to select appropriate document",
    "✓ Error Resilience - Graceful handling, skip on failures",
    "✓ Fully Configurable - No code changes for new matching rules",
    "",
    "Next Steps:",
    "• Add more complex condition types (AND/OR logic)",
    "• Implement caching for API responses",
    "• Add parallel API execution for independent calls"
])

# Save presentation
prs.save('C:/Users/ghmd8/Documents/AI/fresh-doc-hub-poc/Document_Hub_POC_Presentation.pptx')
print("Presentation created: Document_Hub_POC_Presentation.pptx")
